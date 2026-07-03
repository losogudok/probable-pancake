"""Document parsers for PDF, DOCX, XLSX, PPTX files.

Each parser returns (full_text, pages) where pages is a list of per-page text
(for page-level chunk attribution). For non-paginated formats, pages is a
single-element list.
"""
from __future__ import annotations

import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)


def parse_pdf(path: str) -> tuple[str, list[str]]:
    """Parse PDF using PyMuPDF (fitz). Returns (full_text, pages)."""
    import fitz
    doc = fitz.open(path)
    pages = []
    for page in doc:
        text = page.get_text("text")
        pages.append(text)
    doc.close()
    full_text = "\n\n".join(pages)
    return full_text, pages


def parse_docx(path: str) -> tuple[str, list[str]]:
    """Parse DOCX using python-docx."""
    from docx import Document
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(" | ".join(cells))
    full_text = "\n".join(paragraphs)
    return full_text, [full_text]


def parse_doc(path: str) -> tuple[str, list[str]]:
    """Parse old .doc format — try antiword, fallback to docx parser."""
    # Try docx parser first (some .doc files are actually docx)
    try:
        return parse_docx(path)
    except Exception:
        pass
    # Try antiword
    try:
        import subprocess
        result = subprocess.run(
            ["antiword", path], capture_output=True, text=True, timeout=30)
        if result.returncode == 0:
            text = result.stdout
            return text, [text]
    except Exception:
        pass
    logger.warning(f"Could not parse .doc file: {path}")
    return "", []


def parse_xlsx(path: str) -> tuple[str, list[str]]:
    """Parse XLSX/XLS as text representation of all sheets."""
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    pages = []
    for sheet in wb.worksheets:
        lines = []
        for row in sheet.iter_rows(max_row=500, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
        if lines:
            pages.append(f"Sheet: {sheet.title}\n" + "\n".join(lines))
    wb.close()
    return "\n\n".join(pages), pages


def parse_pptx(path: str) -> tuple[str, list[str]]:
    """Parse PPTX presentation."""
    from pptx import Presentation
    prs = Presentation(path)
    pages = []
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        pages.append("\n".join(texts))
    return "\n\n".join(pages), pages


def parse_txt(path: str) -> tuple[str, list[str]]:
    """Parse plain text file."""
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return text, [text]


PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".doc": parse_doc,
    ".docm": parse_docx,
    ".xlsx": parse_xlsx,
    ".xls": parse_xlsx,
    ".pptx": parse_pptx,
    ".txt": parse_txt,
    ".md": parse_txt,
}


def parse_document(path: str) -> tuple[str, list[str]]:
    """Parse a document by extension. Returns (full_text, pages)."""
    ext = os.path.splitext(path)[1].lower()
    parser = PARSERS.get(ext)
    if parser is None:
        logger.warning(f"Unsupported file type: {ext} ({path})")
        return "", []
    try:
        return parser(path)
    except Exception as e:
        logger.error(f"Error parsing {path}: {e}")
        return "", []
